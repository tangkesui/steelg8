"""
项目文件抽取与切块
--------------------

Phase 2 Step 1：.md / .txt（stdlib）
Phase 2 Step 2：+ .docx（python-docx）
Phase 2 Step 3：+ .pdf（pypdf）、.pptx（python-pptx）、.doc（macOS textutil）

切块策略：
- 先按段落（连续空行）粗分
- 按目标 token 数聚合：默认 500（粗估 1 token ≈ 2 中文字 / 4 英文字符）
- 段落本身超长的，再按句号 / 换行二次切
- 相邻 chunk 之间留约 100 字符重叠，减少"信息落在边界被切断"

输出一个 Chunk dataclass 列表，供 embedding 和入库。
"""

from __future__ import annotations

import os
import re
import shutil as _shutil
import subprocess
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any, Iterable, Iterator

from document import chunkers as document_chunkers
from document import registry as document_registry
from document.blocks import ParsedDocument


# ---- 配置常量 ----

# 默认允许的扩展名。某些格式需要 venv 里装了依赖才真正生效
SUPPORTED_EXT: set[str] = {
    ".md", ".markdown", ".txt", ".mdx",
    ".docx",
    ".pdf",
    ".pptx",
    ".doc",
}


def _has_docx() -> bool:
    try:
        import docx  # noqa: F401
        return True
    except ImportError:
        return False


def _has_pypdf() -> bool:
    try:
        import pypdf  # noqa: F401
        return True
    except ImportError:
        return False


def _has_pptx() -> bool:
    try:
        import pptx  # noqa: F401
        return True
    except ImportError:
        return False


def _has_textutil() -> bool:
    """macOS 自带 textutil；非 macOS 或未装 Developer Tools 时返回 False。"""
    return _shutil.which("textutil") is not None

# 目录黑名单（绝对跳过）
SKIP_DIRS: set[str] = {
    ".git", ".svn", ".hg",
    "node_modules", "__pycache__", ".venv", "venv", ".env",
    ".build", ".cache", ".home", ".local-test",
    ".idea", ".vscode", ".claude",
    "dist", "build", ".next", ".nuxt",
    ".steelg8",
}

# 文件大小上限（按类型分档）：二进制格式天生比纯文本大得多
FILE_SIZE_LIMITS: dict[str, int] = {
    ".md":       1_000_000,
    ".markdown": 1_000_000,
    ".mdx":      1_000_000,
    ".txt":      1_000_000,
    ".docx":    50_000_000,
    ".doc":     50_000_000,
    ".pdf":    100_000_000,
    ".pptx":    50_000_000,
}
DEFAULT_SIZE_LIMIT = 1_000_000

# 单个文件抽出来的纯文本最长截到这里，防止一本书式 PDF 把管线爆掉
MAX_EXTRACT_CHARS = 500_000

# 粗估 token 数：每个字约 0.5 token
def _approx_tokens(s: str) -> int:
    return max(1, len(s) // 2)


@dataclass(frozen=True)
class FileRef:
    abs_path: str
    rel_path: str           # 相对 project root
    size: int
    mtime: float


@dataclass
class Chunk:
    rel_path: str
    chunk_idx: int
    text: str
    approx_tokens: int = field(default=0)
    source_path: str = ""
    page: int | None = None
    heading: str = ""
    paragraph_idx: int = 0
    start_char: int = 0
    end_char: int = 0
    content_hash: str = ""
    metadata: dict[str, Any] = field(default_factory=dict)

    def __post_init__(self) -> None:
        if self.approx_tokens == 0:
            object.__setattr__(self, "approx_tokens", _approx_tokens(self.text))
        if not self.source_path:
            self.source_path = self.rel_path
        if not self.end_char:
            self.end_char = self.start_char + len(self.text)
        if not self.content_hash:
            self.content_hash = text_hash(self.text)


@dataclass(frozen=True)
class ParserDiagnostics:
    rel_path: str
    parser: str
    title: str
    block_count: int
    chunk_count: int
    char_count: int
    table_count: int
    code_count: int
    heading_count: int
    list_count: int
    empty_text: bool
    truncated: bool
    chunk_profile: str
    block_types: dict[str, int]

    def to_dict(self) -> dict[str, Any]:
        return {
            "relPath": self.rel_path,
            "parser": self.parser,
            "title": self.title,
            "blockCount": self.block_count,
            "chunkCount": self.chunk_count,
            "charCount": self.char_count,
            "tableCount": self.table_count,
            "codeCount": self.code_count,
            "headingCount": self.heading_count,
            "listCount": self.list_count,
            "emptyText": self.empty_text,
            "truncated": self.truncated,
            "chunkProfile": self.chunk_profile,
            "blockTypes": dict(self.block_types),
        }


# ---- 公开 API ----


def walk_project(root: str) -> Iterator[FileRef]:
    """按 SUPPORTED_EXT 过滤 + SKIP_DIRS 黑名单遍历。按 rel_path 自然排序。"""
    root_path = Path(root).resolve()
    if not root_path.is_dir():
        return

    # 按依赖可用性裁剪：缺包的格式不走索引，免得静默失败
    docx_ok = _has_docx()
    pdf_ok = _has_pypdf()
    pptx_ok = _has_pptx()
    doc_ok = _has_textutil()

    results: list[FileRef] = []
    for dirpath, dirnames, filenames in os.walk(root_path):
        # in-place 修改 dirnames 让 os.walk 不递归黑名单目录
        dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS and not d.startswith(".")]
        for fname in filenames:
            # macOS / Office 临时锁文件
            if fname.startswith("~$") or fname.startswith(".~"):
                continue
            ext = Path(fname).suffix.lower()
            if ext not in SUPPORTED_EXT:
                continue
            if ext == ".docx" and not docx_ok:
                continue
            if ext == ".pdf" and not pdf_ok:
                continue
            if ext == ".pptx" and not pptx_ok:
                continue
            if ext == ".doc" and not doc_ok:
                continue
            abs_path = Path(dirpath) / fname
            try:
                stat = abs_path.stat()
            except OSError:
                continue
            size_limit = FILE_SIZE_LIMITS.get(ext, DEFAULT_SIZE_LIMIT)
            if stat.st_size > size_limit:
                continue
            if stat.st_size == 0:
                continue
            try:
                rel = abs_path.relative_to(root_path)
            except ValueError:
                continue
            results.append(
                FileRef(
                    abs_path=str(abs_path),
                    rel_path=str(rel),
                    size=stat.st_size,
                    mtime=stat.st_mtime,
                )
            )

    results.sort(key=lambda r: r.rel_path)
    yield from results


def parse_document(abs_path: str, *, rel_path: str | None = None) -> ParsedDocument:
    """Parse a source file into structured document blocks.

    This is the new document pipeline entry. `read_text` remains as a legacy
    compatibility wrapper for callers that still expect plain text.
    """
    return document_registry.parse_file(abs_path, rel_path=rel_path, max_chars=MAX_EXTRACT_CHARS)


def read_text(abs_path: str) -> str:
    """把支持的文件读成纯文本。保留旧 API，内部已切到结构化 parser。"""
    return parse_document(abs_path).to_text()


def chunk_document(
    document: ParsedDocument,
    *,
    target_tokens: int = 500,
    overlap_chars: int = 100,
    chunk_profile: str = "default",
) -> list[Chunk]:
    block_chunks = document_chunkers.chunk_document(
        document,
        target_tokens=target_tokens,
        overlap_chars=overlap_chars,
        profile=chunk_profile,
    )
    return [
        Chunk(
            rel_path=c.rel_path,
            chunk_idx=c.chunk_idx,
            text=c.text,
            approx_tokens=c.approx_tokens,
            source_path=c.source_path,
            page=c.page,
            heading=c.heading,
            paragraph_idx=c.paragraph_idx,
            start_char=c.start_char,
            end_char=c.end_char,
            content_hash=c.content_hash,
            metadata=dict(c.metadata or {}),
        )
        for c in block_chunks
    ]


def parser_diagnostics(
    document: ParsedDocument,
    chunks: list[Chunk],
    *,
    chunk_profile: str = "default",
) -> ParserDiagnostics:
    block_types: dict[str, int] = {}
    char_count = 0
    for block in document.blocks:
        block_types[block.type] = block_types.get(block.type, 0) + 1
        char_count += max(0, block.end_char - block.start_char)
    return ParserDiagnostics(
        rel_path=document.rel_path,
        parser=document.parser,
        title=document.title,
        block_count=len(document.blocks),
        chunk_count=len(chunks),
        char_count=char_count,
        table_count=block_types.get("table", 0),
        code_count=block_types.get("code", 0),
        heading_count=block_types.get("heading", 0),
        list_count=block_types.get("list", 0),
        empty_text=not any(block.text.strip() for block in document.blocks),
        truncated=bool(document.metadata.get("truncated")),
        chunk_profile=chunk_profile,
        block_types=block_types,
    )


def file_hash(abs_path: str, *, chunk_size: int = 1024 * 1024) -> str:
    """Return a sha256 hash of the raw file bytes for incremental indexing."""
    import hashlib

    h = hashlib.sha256()
    with open(abs_path, "rb") as f:
        while True:
            b = f.read(chunk_size)
            if not b:
                break
            h.update(b)
    return h.hexdigest()


def text_hash(text: str) -> str:
    import hashlib

    return hashlib.sha256((text or "").encode("utf-8")).hexdigest()


def read_docx(abs_path: str) -> str:
    """用 python-docx 把 docx 抽成纯文本：段落按出现顺序，表格转 Markdown 表。

    未装 python-docx 时返回空字符串（RAG 层会跳过该条）。
    """
    try:
        from docx import Document
    except ImportError:
        return ""

    try:
        doc = Document(abs_path)
    except Exception:
        return ""

    out: list[str] = []

    # body element 有 <w:p> 段落和 <w:tbl> 表格交替出现，保持顺序
    body = doc.element.body
    # XML 本地名映射
    W_NS = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    para_iter = iter(doc.paragraphs)
    table_iter = iter(doc.tables)

    for child in body.iterchildren():
        tag = child.tag
        if tag == f"{{{W_NS}}}p":
            try:
                p = next(para_iter)
            except StopIteration:
                continue
            txt = (p.text or "").strip()
            if not txt:
                continue
            style = getattr(p.style, "name", "") or ""
            # 标题带 markdown 井号，便于下游 chunker 识别
            if style.startswith("Heading 1"):
                out.append(f"# {txt}")
            elif style.startswith("Heading 2"):
                out.append(f"## {txt}")
            elif style.startswith("Heading 3"):
                out.append(f"### {txt}")
            elif style.startswith("Heading 4"):
                out.append(f"#### {txt}")
            else:
                out.append(txt)
        elif tag == f"{{{W_NS}}}tbl":
            try:
                tbl = next(table_iter)
            except StopIteration:
                continue
            md = _docx_table_to_markdown(tbl)
            if md:
                out.append(md)

    return "\n\n".join(out)


def read_pdf(abs_path: str) -> str:
    """用 pypdf 抽文本 PDF 的文字。扫描件（图像）得靠 OCR，不在这里管。"""
    try:
        import pypdf
    except ImportError:
        return ""

    try:
        reader = pypdf.PdfReader(abs_path)
    except Exception:
        return ""

    pages_text: list[str] = []
    for i, page in enumerate(reader.pages, start=1):
        try:
            txt = page.extract_text() or ""
        except Exception:
            txt = ""
        txt = txt.strip()
        if txt:
            pages_text.append(f"<!-- page {i} -->\n{txt}")

    return "\n\n".join(pages_text)


def read_pptx(abs_path: str) -> str:
    """从 pptx 抽文本：每页一个二级标题，每个 text frame 一段。"""
    try:
        from pptx import Presentation
    except ImportError:
        return ""

    try:
        prs = Presentation(abs_path)
    except Exception:
        return ""

    out: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        # 标题通常在第一个 placeholder
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            full = "\n".join(p.text for p in text_frame.paragraphs if p.text).strip()
            if not full:
                continue
            # 取第一个非空的作为标题
            if not title:
                title = full.split("\n")[0][:80]
                rest = "\n".join(full.split("\n")[1:]).strip()
                if rest:
                    body_parts.append(rest)
            else:
                body_parts.append(full)

        if title:
            out.append(f"## {title}")
        else:
            out.append(f"## Slide {idx}")
        if body_parts:
            out.append("\n\n".join(body_parts))

    return "\n\n".join(out)


def read_legacy_doc(abs_path: str) -> str:
    """用 macOS 自带的 textutil 把 .doc 转 txt，再读。
    非 macOS / 未装开发者工具时返回空。"""
    if not _has_textutil():
        return ""
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        tmp_path = f.name
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-output", tmp_path, abs_path],
            capture_output=True, timeout=30,
        )
        if result.returncode != 0:
            return ""
        try:
            with open(tmp_path, "r", encoding="utf-8") as f:
                return f.read()
        except UnicodeDecodeError:
            with open(tmp_path, "r", encoding="gbk", errors="replace") as f:
                return f.read()
    except (subprocess.SubprocessError, OSError):
        return ""
    finally:
        try:
            os.unlink(tmp_path)
        except OSError:
            pass


def _docx_table_to_markdown(tbl: object) -> str:
    """转成 Markdown pipe table。首行当表头；若表头和后续行列数不齐，按最长列数补空。"""
    rows: list[list[str]] = []
    for row in getattr(tbl, "rows", []):
        cells = [(c.text or "").strip().replace("\n", " ").replace("|", "\\|") for c in row.cells]
        rows.append(cells)
    if not rows:
        return ""
    # 去掉完全空的行
    rows = [r for r in rows if any(cell for cell in r)]
    if not rows:
        return ""

    max_cols = max(len(r) for r in rows)
    rows = [r + [""] * (max_cols - len(r)) for r in rows]

    header = rows[0]
    body = rows[1:]
    md_lines = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * max_cols) + " |",
    ]
    for r in body:
        md_lines.append("| " + " | ".join(r) + " |")
    return "\n".join(md_lines)


def chunk_text(
    text: str,
    rel_path: str,
    *,
    target_tokens: int = 500,
    overlap_chars: int = 100,
    chunk_profile: str = "default",
) -> list[Chunk]:
    """把单个文件的全文切成 Chunk 列表。"""
    if not text.strip():
        return []

    document = document_registry.parse_text(text, rel_path=rel_path, parser="legacy-text")
    return chunk_document(
        document,
        target_tokens=target_tokens,
        overlap_chars=overlap_chars,
        chunk_profile=chunk_profile,
    )


@dataclass
class _Paragraph:
    text: str
    start_char: int
    end_char: int
    paragraph_idx: int
    page: int | None
    heading: str


_PAGE_MARKER_RE = re.compile(r"<!--\s*page\s+(\d+)\s*-->", re.IGNORECASE)
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")


def _paragraphs_with_metadata(text: str) -> list[_Paragraph]:
    normalized = text.replace("\r\n", "\n")
    out: list[_Paragraph] = []
    page: int | None = None
    heading = ""
    paragraph_idx = 0

    for match in re.finditer(r"\S.*?(?=\n\s*\n|\Z)", normalized, flags=re.DOTALL):
        raw = match.group(0)
        leading = len(raw) - len(raw.lstrip())
        trailing = len(raw.rstrip())
        start = match.start() + leading
        end = match.start() + trailing
        para = raw.strip()
        if not para:
            continue

        page_match = _PAGE_MARKER_RE.search(para)
        if page_match:
            try:
                page = int(page_match.group(1))
            except ValueError:
                pass
            para = _PAGE_MARKER_RE.sub("", para).strip()
            if not para:
                continue

        heading_match = _HEADING_RE.match(para.splitlines()[0].strip())
        if heading_match:
            heading = heading_match.group(2).strip()

        out.append(
            _Paragraph(
                text=para,
                start_char=start,
                end_char=end,
                paragraph_idx=paragraph_idx,
                page=page,
                heading=heading,
            )
        )
        paragraph_idx += 1
    return out


def _chunk_from_parts(
    rel_path: str,
    chunk_idx: int,
    body: str,
    parts: list[_Paragraph],
) -> Chunk:
    first = parts[0]
    last = parts[-1]
    return Chunk(
        rel_path=rel_path,
        chunk_idx=chunk_idx,
        text=body,
        source_path=rel_path,
        page=first.page,
        heading=first.heading,
        paragraph_idx=first.paragraph_idx,
        start_char=first.start_char,
        end_char=last.end_char,
    )


def _split_long_paragraph(para: str, target_tokens: int) -> Iterable[str]:
    """按中英文句尾分隔符切长段落。"""
    sentinels = ["。", "！", "？", ".", "!", "?", "\n"]
    # 简单按字符扫描，在 sentinel 后断开
    pieces: list[str] = []
    current = []
    cur_tokens = 0
    for ch in para:
        current.append(ch)
        cur_tokens += 1 if not ch.isspace() else 0
        if ch in sentinels and cur_tokens >= target_tokens * 2:
            pieces.append("".join(current).strip())
            current = []
            cur_tokens = 0
    if current:
        tail = "".join(current).strip()
        if tail:
            pieces.append(tail)
    return pieces


def _split_long_paragraph_with_spans(para: _Paragraph, target_tokens: int) -> Iterable[_Paragraph]:
    pieces = list(_split_long_paragraph(para.text, target_tokens))
    cursor = 0
    for idx, piece in enumerate(pieces):
        rel_start = para.text.find(piece, cursor)
        if rel_start < 0:
            rel_start = cursor
        rel_end = rel_start + len(piece)
        cursor = rel_end
        yield _Paragraph(
            text=piece,
            start_char=para.start_char + rel_start,
            end_char=para.start_char + rel_end,
            paragraph_idx=para.paragraph_idx + idx,
            page=para.page,
            heading=para.heading,
        )


# ---- 一把梭便捷函数 ----


def extract_and_chunk(root: str, *, target_tokens: int = 500) -> list[Chunk]:
    """遍历 root 下所有支持的文件，切成 chunk 列表。"""
    out: list[Chunk] = []
    for fref in walk_project(root):
        text = read_text(fref.abs_path)
        chunks = chunk_text(text, fref.rel_path, target_tokens=target_tokens)
        out.extend(chunks)
    return out
