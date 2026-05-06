from __future__ import annotations

import hashlib
import re
import shutil
import subprocess
from pathlib import Path
from typing import Any

from .blocks import DocumentBlock, ParsedDocument, fallback_title


DEFAULT_MAX_CHARS = 500_000

_PAGE_MARKER_RE = re.compile(r"^\s*<!--\s*page\s+(\d+)\s*-->\s*$", re.IGNORECASE)
_HEADING_RE = re.compile(r"^(#{1,6})\s+(.+?)\s*$")
_TABLE_ROW_RE = re.compile(r"^\s*\|.*\|\s*$")
_LIST_RE = re.compile(r"^\s*(?:[-*+]|\d+[.)])\s+")
_CODE_FENCE_RE = re.compile(r"^\s*(```|~~~)")


def parse_file(abs_path: str, *, rel_path: str | None = None, max_chars: int = DEFAULT_MAX_CHARS) -> ParsedDocument:
    path = Path(abs_path)
    rel = rel_path or path.name
    ext = path.suffix.lower()
    if ext == ".docx":
        text = read_docx_text(abs_path)
        parser = "docx"
    elif ext == ".pdf":
        text = read_pdf_text(abs_path)
        parser = "pdf"
    elif ext == ".pptx":
        text = read_pptx_text(abs_path)
        parser = "pptx"
    elif ext == ".doc":
        text = read_legacy_doc_text(abs_path)
        parser = "doc"
    else:
        text = read_plain_text(abs_path)
        parser = "markdown" if ext in {".md", ".markdown", ".mdx"} else "text"
    raw_chars = len(text)
    truncated = raw_chars > max_chars
    if truncated:
        text = text[:max_chars] + "\n\n…（抽取内容已截断，超过 {:,} 字）".format(max_chars)
    return parse_text(
        text,
        rel_path=rel,
        parser=parser,
        metadata={
            "parser": parser,
            "raw_chars": raw_chars,
            "truncated": truncated,
            "max_chars": max_chars,
        },
    )


def parse_text(
    text: str,
    *,
    rel_path: str,
    parser: str = "text",
    metadata: dict[str, Any] | None = None,
) -> ParsedDocument:
    normalized = (text or "").replace("\r\n", "\n").replace("\r", "\n")
    blocks = _blocks_from_text(normalized, parser=parser)
    title = _title_from_blocks(blocks) or fallback_title(rel_path)
    doc_metadata = {"parser": parser}
    doc_metadata.update(metadata or {})
    return ParsedDocument(
        rel_path=rel_path,
        title=title,
        blocks=blocks,
        parser=parser,
        content_hash=_hash_text(normalized),
        metadata=doc_metadata,
    )


def read_plain_text(abs_path: str) -> str:
    try:
        return Path(abs_path).read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return Path(abs_path).read_text(encoding="latin-1")


def read_docx_text(abs_path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""
    try:
        doc = Document(abs_path)
    except Exception:
        return ""

    out: list[str] = []
    body = doc.element.body
    para_iter = iter(doc.paragraphs)
    table_iter = iter(doc.tables)
    w_ns = "http://schemas.openxmlformats.org/wordprocessingml/2006/main"

    for child in body.iterchildren():
        tag = child.tag
        if tag == f"{{{w_ns}}}p":
            try:
                paragraph = next(para_iter)
            except StopIteration:
                continue
            text = (paragraph.text or "").strip()
            if not text:
                continue
            style = getattr(paragraph.style, "name", "") or ""
            if style.startswith("Heading 1"):
                out.append(f"# {text}")
            elif style.startswith("Heading 2"):
                out.append(f"## {text}")
            elif style.startswith("Heading 3"):
                out.append(f"### {text}")
            elif style.startswith("Heading 4"):
                out.append(f"#### {text}")
            else:
                out.append(text)
        elif tag == f"{{{w_ns}}}tbl":
            try:
                table = next(table_iter)
            except StopIteration:
                continue
            table_text = _docx_table_to_markdown(table)
            if table_text:
                out.append(table_text)
    return "\n\n".join(out)


def read_pdf_text(abs_path: str) -> str:
    try:
        import pypdf
    except ImportError:
        return ""
    try:
        reader = pypdf.PdfReader(abs_path)
    except Exception:
        return ""

    pages: list[str] = []
    for idx, page in enumerate(reader.pages, start=1):
        try:
            text = page.extract_text() or ""
        except Exception:
            text = ""
        text = text.strip()
        if text:
            pages.append(f"<!-- page {idx} -->\n{text}")
    return "\n\n".join(pages)


def read_pptx_text(abs_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        presentation = Presentation(abs_path)
    except Exception:
        return ""

    out: list[str] = []
    for idx, slide in enumerate(presentation.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text_frame = shape.text_frame
            full = "\n".join(p.text for p in text_frame.paragraphs if p.text).strip()
            if not full:
                continue
            if not title:
                title = full.split("\n")[0][:80]
                rest = "\n".join(full.split("\n")[1:]).strip()
                if rest:
                    body_parts.append(rest)
            else:
                body_parts.append(full)

        out.append(f"<!-- page {idx} -->")
        out.append(f"## {title or f'Slide {idx}'}")
        if body_parts:
            out.append("\n\n".join(body_parts))
    return "\n\n".join(out)


def read_legacy_doc_text(abs_path: str) -> str:
    if shutil.which("textutil") is None:
        return ""
    import tempfile

    with tempfile.NamedTemporaryFile(suffix=".txt", delete=False) as f:
        tmp_path = f.name
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-output", tmp_path, abs_path],
            capture_output=True,
            timeout=30,
        )
        if result.returncode != 0:
            return ""
        try:
            return Path(tmp_path).read_text(encoding="utf-8")
        except UnicodeDecodeError:
            return Path(tmp_path).read_text(encoding="gbk", errors="replace")
    except (OSError, subprocess.SubprocessError):
        return ""
    finally:
        try:
            Path(tmp_path).unlink()
        except OSError:
            pass


def _blocks_from_text(text: str, *, parser: str) -> list[DocumentBlock]:
    lines = text.splitlines(keepends=True)
    offsets: list[int] = []
    cursor = 0
    for line in lines:
        offsets.append(cursor)
        cursor += len(line)

    blocks: list[DocumentBlock] = []
    heading_stack: list[str] = []
    page: int | None = None
    i = 0
    while i < len(lines):
        raw_line = lines[i]
        line = raw_line.strip()
        if not line:
            i += 1
            continue

        page_match = _PAGE_MARKER_RE.match(line)
        if page_match:
            try:
                page = int(page_match.group(1))
            except ValueError:
                pass
            i += 1
            continue

        code_match = _CODE_FENCE_RE.match(line)
        if code_match:
            fence = code_match.group(1)
            j = i + 1
            while j < len(lines):
                if lines[j].strip().startswith(fence):
                    j += 1
                    break
                j += 1
            blocks.append(_make_block("code", lines, offsets, i, j, page, heading_stack, parser))
            i = j
            continue

        heading_match = _HEADING_RE.match(line)
        if heading_match:
            level = len(heading_match.group(1))
            title = heading_match.group(2).strip()
            heading_stack = heading_stack[: max(0, level - 1)] + [title]
            blocks.append(
                _make_block(
                    "heading",
                    lines,
                    offsets,
                    i,
                    i + 1,
                    page,
                    heading_stack,
                    parser,
                    extra={"level": level, "title": title},
                )
            )
            i += 1
            continue

        if _TABLE_ROW_RE.match(line):
            j = i + 1
            while j < len(lines) and _TABLE_ROW_RE.match(lines[j].strip()):
                j += 1
            blocks.append(_make_block("table", lines, offsets, i, j, page, heading_stack, parser))
            i = j
            continue

        block_type = "list" if _LIST_RE.match(line) else "paragraph"
        j = i + 1
        while j < len(lines):
            probe = lines[j].strip()
            if not probe:
                break
            if _PAGE_MARKER_RE.match(probe) or _HEADING_RE.match(probe) or _CODE_FENCE_RE.match(probe):
                break
            if block_type != "table" and _TABLE_ROW_RE.match(probe):
                break
            if block_type == "list" and not _LIST_RE.match(probe):
                break
            j += 1
        blocks.append(_make_block(block_type, lines, offsets, i, j, page, heading_stack, parser))
        i = j
    for idx, block in enumerate(blocks):
        block.metadata.setdefault("block_index", idx)
        block.metadata.setdefault("paragraph_idx", idx)
    return blocks


def _make_block(
    block_type: str,
    lines: list[str],
    offsets: list[int],
    start_line: int,
    end_line: int,
    page: int | None,
    heading_stack: list[str],
    parser: str,
    *,
    extra: dict | None = None,
) -> DocumentBlock:
    raw = "".join(lines[start_line:end_line])
    leading = len(raw) - len(raw.lstrip())
    trailing_text = raw.rstrip()
    start = offsets[start_line] + leading
    end = offsets[start_line] + len(trailing_text)
    metadata = {"parser": parser}
    if extra:
        metadata.update(extra)
    return DocumentBlock(
        type=block_type,
        text=raw.strip(),
        start_char=start,
        end_char=end,
        page=page,
        heading_path=list(heading_stack),
        metadata=metadata,
    )


def _title_from_blocks(blocks: list[DocumentBlock]) -> str:
    for block in blocks:
        if block.type == "heading" and block.heading:
            return block.heading
    return ""


def _docx_table_to_markdown(table: object) -> str:
    rows: list[list[str]] = []
    for row in getattr(table, "rows", []):
        cells = [(c.text or "").strip().replace("\n", " ").replace("|", "\\|") for c in row.cells]
        rows.append(cells)
    rows = [row for row in rows if any(cell for cell in row)]
    if not rows:
        return ""
    max_cols = max(len(row) for row in rows)
    rows = [row + [""] * (max_cols - len(row)) for row in rows]
    header = rows[0]
    body = rows[1:]
    out = [
        "| " + " | ".join(header) + " |",
        "| " + " | ".join(["---"] * max_cols) + " |",
    ]
    for row in body:
        out.append("| " + " | ".join(row) + " |")
    return "\n".join(out)


def _hash_text(text: str) -> str:
    return hashlib.sha256((text or "").encode("utf-8")).hexdigest()
